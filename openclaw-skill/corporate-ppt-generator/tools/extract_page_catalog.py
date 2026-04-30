from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def guess_page_type(layout_name: str, texts: list[str]) -> str:
    name = layout_name.lower()
    merged = " ".join(texts).lower()
    if "cover" in name or "封面" in name:
        return "cover"
    if "目录" in name or "content" in name:
        return "toc"
    if "章节页" in name or "section page" in name:
        return "section"
    if "副章节" in name or "sub section" in name:
        return "sub_section"
    if "end page" in name or "封底" in name or "thanks" in merged:
        return "end"
    if any(k in merged for k in ["饼状图", "柱状图", "折线图", "面积图", "图表", "chart", "table", "地图"]):
        return "chart_or_data"
    if "标准内容页" in name or "standard page" in name:
        return "content"
    return "generic"


def extract_catalog(template_path: Path, output_path: Path) -> dict:
    prs = Presentation(str(template_path))
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        placeholders = []
        shape_stats = {
            "total": len(slide.shapes),
            "placeholder": 0,
            "picture": 0,
            "table": 0,
            "chart": 0,
            "group": 0,
            "textbox_or_text": 0,
        }

        for shape in slide.shapes:
            st = str(getattr(shape, "shape_type", ""))
            is_placeholder = bool(getattr(shape, "is_placeholder", False))
            has_text_frame = bool(getattr(shape, "has_text_frame", False))
            has_table = bool(getattr(shape, "has_table", False))
            has_chart = bool(getattr(shape, "has_chart", False))

            if is_placeholder:
                shape_stats["placeholder"] += 1
                if has_text_frame:
                    placeholders.append(
                        {
                            "idx": shape.placeholder_format.idx,
                            "name": getattr(shape, "name", ""),
                        }
                    )
            if "PICTURE" in st:
                shape_stats["picture"] += 1
            if has_table:
                shape_stats["table"] += 1
            if has_chart:
                shape_stats["chart"] += 1
            if "GROUP" in st:
                shape_stats["group"] += 1
            if has_text_frame:
                shape_stats["textbox_or_text"] += 1
                txt = shape.text.strip().replace("\n", " / ")
                if txt:
                    texts.append(txt[:140])

        placeholders.sort(key=lambda item: item["idx"])
        page_type = guess_page_type(slide.slide_layout.name, texts)
        pages.append(
            {
                "index": i,
                "layout_name": slide.slide_layout.name,
                "page_type_guess": page_type,
                "sample_texts": texts[:5],
                "placeholder_schema": placeholders,
                "shape_stats": shape_stats,
            }
        )

    result = {
        "template_file": template_path.name,
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "pages": pages,
    }
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    return result


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract page catalog from PPT template.")
    parser.add_argument("--template", required=True, help="Path to template .pptx")
    parser.add_argument("--output", required=True, help="Path to output page_catalog.json")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    template_path = Path(args.template)
    output_path = Path(args.output)
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    result = extract_catalog(template_path, output_path)
    print(f"Catalog generated: {output_path.resolve()}")
    print(f"Slides archived: {result['slide_count']}")


if __name__ == "__main__":
    main()
