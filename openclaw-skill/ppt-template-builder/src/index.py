from __future__ import annotations

import argparse
import asyncio
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _find_layout_by_name(prs: Presentation, keyword: str):
    for layout in prs.slide_layouts:
        if keyword in layout.name:
            return layout
    raise ValueError(f"Layout not found with keyword: {keyword}")


def _fill_placeholders(slide, texts: list[str]) -> None:
    placeholders = []
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        placeholders.append((shape.placeholder_format.idx, shape))
    placeholders.sort(key=lambda item: item[0])
    for i, text in enumerate(texts):
        if i >= len(placeholders):
            break
        placeholders[i][1].text = text


def _fill_placeholder_map(slide, text_by_idx: dict[int, str]) -> None:
    by_idx = {}
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        by_idx[shape.placeholder_format.idx] = shape
    for idx, text in text_by_idx.items():
        target = by_idx.get(idx)
        if target is not None:
            target.text = text


def _fill_toc_slide(slide, title: str, items: list[str]) -> None:
    _fill_placeholder_map(
        slide,
        {
            0: title or "目录",
            11: "Content",
            12: "\n".join(items),
        },
    )


def _set_nonempty_texts(slide, texts: list[str]) -> None:
    targets = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.strip():
            targets.append(shape)
    for i, text in enumerate(texts):
        if i >= len(targets):
            break
        targets[i].text = text


def _iter_text_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text_shapes(shape.shapes)
            continue
        if getattr(shape, "has_text_frame", False):
            yield shape


def _fill_group_cards(slide, headings: list[str], bodies: list[str]) -> None:
    heading_shapes = []
    body_shapes = []
    for shape in _iter_text_shapes(slide.shapes):
        text = (shape.text or "").strip().lower()
        if "text here" in text:
            heading_shapes.append(shape)
            continue
        if "lorem ipsum" in text:
            body_shapes.append(shape)
    for i, text in enumerate(headings):
        if i >= len(heading_shapes):
            break
        heading_shapes[i].text = text
    for i, text in enumerate(bodies):
        if i >= len(body_shapes):
            break
        body_shapes[i].text = text


def _retain_slide_indices(prs: Presentation, keep_indices: set[int]) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for idx in range(len(prs.slides) - 1, -1, -1):
        if idx in keep_indices:
            continue
        sld_id = sld_id_lst[idx]
        prs.part.drop_rel(sld_id.rId)
        del sld_id_lst[idx]


def _reorder_retained_slides(prs: Presentation, template_order: list[int], selected_sorted: list[int]) -> None:
    if not template_order:
        return
    sld_id_lst = prs.slides._sldIdLst
    template_to_elem = {selected_sorted[i]: sld_id_lst[i] for i in range(len(selected_sorted))}
    original = list(sld_id_lst)
    for elem in original:
        sld_id_lst.remove(elem)
    for template_idx in template_order:
        elem = template_to_elem.get(template_idx)
        if elem is not None:
            sld_id_lst.append(elem)


def _update_first_chart(slide, categories: list[str], series_name: str, values: list[float]) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        data = ChartData()
        data.categories = categories
        data.add_series(series_name, values)
        shape.chart.replace_data(data)
        return


def _update_first_table(slide, headers: list[str], rows: list[list[str]]) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        col_count = min(len(headers), len(table.columns))
        for c in range(col_count):
            table.cell(0, c).text = headers[c]
        row_count = min(len(rows), len(table.rows) - 1)
        for r in range(row_count):
            for c in range(min(len(rows[r]), len(table.columns))):
                table.cell(r + 1, c).text = str(rows[r][c])
        return


def _child_text(node: ET.Element, tag_name: str) -> str:
    child = node.find(tag_name)
    if child is None:
        return ""
    return "".join(child.itertext()).strip()


def _split_items(raw: str) -> list[str]:
    if not raw:
        return []
    out = []
    for part in raw.replace("\n", ";").split(";"):
        token = part.strip()
        if token:
            out.append(token)
    return out


def _split_csv(raw: str) -> list[str]:
    if not raw:
        return []
    normalized = raw.replace("，", ",").replace("；", ",").replace(";", ",")
    return [token.strip() for token in normalized.split(",") if token.strip()]


def _parse_numbers(raw: str) -> list[float]:
    values: list[float] = []
    for token in _split_csv(raw):
        try:
            values.append(float(token))
        except ValueError:
            continue
    return values


def _parse_rows(node: ET.Element) -> list[list[str]]:
    rows: list[list[str]] = []
    rows_node = node.find("rows")
    if rows_node is None:
        return rows
    for row in rows_node.findall("row"):
        cells = ["".join(cell.itertext()).strip() for cell in row.findall("cell")]
        if not cells:
            cells = _split_csv("".join(row.itertext()).strip())
        if any(cells):
            rows.append(cells)
    return rows


def _strip_html_tags(raw_html: str) -> str:
    if not raw_html:
        return ""
    text = re.sub(r"<\s*br\s*/?\s*>", "\n", raw_html, flags=re.I)
    text = re.sub(r"</\s*p\s*>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _parse_block_xml(block_xml: str) -> list[dict[str, Any]]:
    root = ET.fromstring(block_xml)
    slide_nodes = [root] if root.tag.lower() == "slide" else [node for node in root.findall(".//slide")]
    blocks: list[dict[str, Any]] = []
    for node in slide_nodes:
        block = {
            "type": (node.attrib.get("type") or "content").strip().lower(),
            "title": node.attrib.get("title", "").strip() or _child_text(node, "title"),
            "subtitle": node.attrib.get("subtitle", "").strip() or _child_text(node, "subtitle"),
            "body": node.attrib.get("body", "").strip() or _child_text(node, "body"),
            "part": node.attrib.get("part", "").strip() or _child_text(node, "part"),
            "part_no": node.attrib.get("part_no", "").strip() or _child_text(node, "part_no"),
            "html": _child_text(node, "html"),
            "series": node.attrib.get("series", "").strip() or _child_text(node, "series"),
            "categories": _split_csv(node.attrib.get("categories", "").strip() or _child_text(node, "categories")),
            "values": _parse_numbers(node.attrib.get("values", "").strip() or _child_text(node, "values")),
            "headers": _split_csv(node.attrib.get("headers", "").strip() or _child_text(node, "headers")),
            "rows": _parse_rows(node),
        }
        item_nodes = node.findall("./items/item")
        block["items"] = (
            ["".join(it.itertext()).strip() for it in item_nodes if "".join(it.itertext()).strip()]
            if item_nodes
            else _split_items(node.attrib.get("items", "").strip() or _child_text(node, "items"))
        )
        heading_nodes = node.findall("./headings/item")
        block["headings"] = (
            ["".join(it.itertext()).strip() for it in heading_nodes if "".join(it.itertext()).strip()]
            if heading_nodes
            else _split_items(node.attrib.get("headings", "").strip() or _child_text(node, "headings"))
        )
        body_nodes = node.findall("./bodies/item")
        block["bodies"] = (
            ["".join(it.itertext()).strip() for it in body_nodes if "".join(it.itertext()).strip()]
            if body_nodes
            else _split_items(node.attrib.get("bodies", "").strip() or _child_text(node, "bodies"))
        )
        blocks.append(block)
    return blocks


def _allocate_template_pages(
    blocks: list[dict[str, Any]], render_strategy: str, allow_html_fallback: bool
) -> tuple[list[int | None], list[int]]:
    pools = {
        "cover": [0],
        "toc": [1],
        "section": [2, 4, 16, 20, 40, 43],
        "market_section": [20],
        "content": [3, 6, 8, 9, 11, 12],
        "matrix": [15],
        "pie": [21, 22],
        "bar": [23, 24, 25],
        "line": [27, 26],
        "table": [28, 29, 30],
        "slogan": [41, 42],
        "end": [55],
        "html": [3, 6, 8, 9],
        "richtext": [3, 6, 8, 9],
    }
    used: set[int] = set()
    allocations: list[int | None] = []
    for block in blocks:
        chosen = None
        for candidate in pools.get(block.get("type", "content"), []):
            if candidate not in used:
                chosen = candidate
                break
        if chosen is not None:
            used.add(chosen)
            allocations.append(chosen)
            continue
        if render_strategy == "template_only" or not allow_html_fallback:
            allocations.append(None)
            continue
        for candidate in pools["content"]:
            if candidate not in used:
                used.add(candidate)
                chosen = candidate
                break
        allocations.append(chosen)
    template_order: list[int] = []
    for alloc in allocations:
        if alloc is None or alloc in template_order:
            continue
        template_order.append(alloc)
    return allocations, template_order


def _build_from_block_xml(
    template_path: Path,
    output_path: Path,
    deck_title: str,
    block_xml: str,
    render_strategy: str,
    allow_html_fallback: bool,
) -> tuple[int, int]:
    if not block_xml.strip():
        raise ValueError("block_xml cannot be empty")
    blocks = _parse_block_xml(block_xml)
    if not blocks:
        raise ValueError("No <slide> nodes found in block_xml")

    allocations, template_order = _allocate_template_pages(blocks, render_strategy, allow_html_fallback)
    selected = sorted({idx for idx in allocations if idx is not None})
    prs = Presentation(str(template_path))
    if selected:
        _retain_slide_indices(prs, set(selected))
        _reorder_retained_slides(prs, template_order, selected)
    else:
        raise ValueError("No template pages allocated for xml blocks")

    allocated_slides = {template_order[i]: prs.slides[i] for i in range(min(len(template_order), len(prs.slides)))}
    fallback_count = 0
    for block_i, block in enumerate(blocks):
        slide = allocated_slides.get(allocations[block_i]) if allocations[block_i] is not None else None
        block_type = block.get("type", "content")
        if slide is None:
            if render_strategy != "template_only" and allow_html_fallback:
                slide = prs.slides.add_slide(_find_layout_by_name(prs, "标准内容页（小标题）_Standard page with subtitle"))
                compact = block.get("body") or block.get("title") or json.dumps(block, ensure_ascii=False)
                _fill_placeholders(slide, [f"Fallback - {block_type}", compact, "Auto fallback"])
                fallback_count += 1
            continue

        if block_type == "cover":
            _fill_placeholder_map(
                slide,
                {
                    0: block.get("title") or deck_title or "Corporate Deck",
                    1: block.get("subtitle") or "Generated by template-first XML rendering",
                },
            )
        elif block_type == "toc":
            _fill_toc_slide(slide, block.get("title") or "目录", block.get("items") or [])
        elif block_type in ("section", "market_section"):
            _fill_placeholder_map(
                slide,
                {
                    0: block.get("title") or f"Section {block_i + 1}",
                    1: block.get("subtitle") or "Section summary",
                    10: block.get("part") or "Part",
                    11: block.get("part_no") or f"{block_i + 1:02d}",
                },
            )
        elif block_type in ("content", "html", "richtext"):
            body = block.get("body") or ""
            if block.get("html"):
                body = _strip_html_tags(block["html"])
                fallback_count += 1
            _fill_placeholder_map(
                slide,
                {
                    0: block.get("title") or f"Content {block_i + 1}",
                    1: body or "TBD",
                    14: block.get("subtitle") or "Subtitle",
                },
            )
        elif block_type == "matrix":
            _fill_placeholder_map(slide, {0: block.get("title") or "Core Capability Matrix"})
            _fill_group_cards(slide, block.get("headings") or [], block.get("bodies") or [])
        elif block_type in ("pie", "bar", "line"):
            _set_nonempty_texts(slide, [block.get("title") or "", block.get("subtitle") or ""])
            _update_first_chart(
                slide,
                block.get("categories") or ["A", "B", "C"],
                block.get("series") or "Series",
                block.get("values") or [30, 40, 30],
            )
        elif block_type == "table":
            _set_nonempty_texts(slide, [block.get("title") or "Table"])
            _update_first_table(
                slide,
                block.get("headers") or ["Col1", "Col2", "Col3"],
                block.get("rows") or [["A", "B", "C"], ["D", "E", "F"]],
            )
        elif block_type == "slogan":
            _set_nonempty_texts(slide, [block.get("title") or "Keyword", block.get("subtitle") or "Slide"])
        elif block_type == "end":
            _fill_placeholder_map(slide, {13: block.get("title") or "Thanks.", 14: block.get("subtitle") or ""})
        elif render_strategy != "template_only" and allow_html_fallback:
            extra = prs.slides.add_slide(_find_layout_by_name(prs, "标准内容页（小标题）_Standard page with subtitle"))
            compact = block.get("body") or block.get("title") or json.dumps(block, ensure_ascii=False)
            _fill_placeholders(extra, [f"Fallback - {block_type}", compact, "Auto fallback"])
            fallback_count += 1

    prs.save(str(output_path))
    return len(prs.slides), fallback_count


async def handler(input: dict[str, Any], _context: Any) -> dict[str, Any]:
    skill_root = Path(__file__).resolve().parent.parent
    template_file = input.get("template_file") or "PPT_Template.pptx"
    template_path = skill_root / template_file
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    mode = (input.get("mode") or "xml").strip().lower()
    if mode != "xml":
        raise ValueError("Only xml mode is supported")

    title = input.get("title", "Corporate Deck")
    output_filename = input.get("output_filename", "openclaw_generated_xml.pptx")
    output_path = skill_root / output_filename
    render_strategy = input.get("render_strategy", "template_first")
    allow_html_fallback = bool(input.get("allow_html_fallback", True))
    block_xml = input.get("block_xml", "")

    slide_count, fallback_count = _build_from_block_xml(
        template_path=template_path,
        output_path=output_path,
        deck_title=title,
        block_xml=block_xml,
        render_strategy=render_strategy,
        allow_html_fallback=allow_html_fallback,
    )
    return {
        "output_path": str(output_path.resolve()),
        "slide_count": slide_count,
        "fallback_count": fallback_count,
        "message": "XML deck generated with corporate template reuse",
    }


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run ppt-template-builder in xml mode.")
    parser.add_argument("--title", default="Corporate Deck")
    parser.add_argument("--output", default="openclaw_generated_xml.pptx")
    parser.add_argument("--block-xml-file", required=True)
    parser.add_argument("--render-strategy", default="template_first", choices=["template_first", "template_only"])
    parser.add_argument("--allow-html-fallback", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = _parse_args()
    block_xml = Path(args.block_xml_file).read_text(encoding="utf-8")
    payload = {
        "mode": "xml",
        "title": args.title,
        "output_filename": args.output,
        "block_xml": block_xml,
        "render_strategy": args.render_strategy,
        "allow_html_fallback": bool(args.allow_html_fallback),
    }
    result = asyncio.run(handler(payload, None))
    print(result)


if __name__ == "__main__":
    main()
