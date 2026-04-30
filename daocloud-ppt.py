from pathlib import Path
import json
import zipfile

from pptx import Presentation


TEMPLATE_PATH = Path("PPT_Template.pptx")
ASSET_DIR = Path("extracted_company_assets")
OUTPUT_PATH = Path("cc.pptx")
KB_PATH = Path("template_knowledge_base.json")


def extract_media(template_path: Path, output_dir: Path) -> list[Path]:
    """Extract media files from a template pptx package."""
    output_dir.mkdir(parents=True, exist_ok=True)
    extracted: list[Path] = []
    with zipfile.ZipFile(template_path, "r") as zf:
        media_files = sorted(
            name for name in zf.namelist() if name.startswith("ppt/media/") and not name.endswith("/")
        )
        for media_name in media_files:
            data = zf.read(media_name)
            filename = Path(media_name).name
            dst = output_dir / filename
            dst.write_bytes(data)
            extracted.append(dst)
    return extracted


def build_template_knowledge_base(template_path: Path, output_path: Path) -> dict:
    """Build a compact knowledge base from template layouts and sample slides."""
    prs = Presentation(str(template_path))
    layouts = [{"index": i, "name": layout.name} for i, layout in enumerate(prs.slide_layouts)]

    sample_slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                txt = shape.text.strip().replace("\n", " / ")
                if txt:
                    texts.append(txt[:120])
        sample_slides.append(
            {
                "index": i,
                "layout_name": slide.slide_layout.name,
                "sample_texts": texts[:3],
            }
        )

    semantic_layout_map = {
        "cover": "封面_Cover",
        "toc": "目录页_Content",
        "section": "章节页_Section page",
        "sub_section": "副章节页_Sub section page",
        "content": "标准内容页_Standard page",
        "content_with_subtitle": "标准内容页（小标题）_Standard page with subtitle",
        "title_only_bg": "仅标题（带背景）_Title only with bg",
        "title_only_plain": "仅标题（无背景）_Title only without bg",
        "slogan_dark": "标语_Slogan_Dark",
        "slogan_light": "标语_Slogan_Light",
        "end": "封底_End page",
    }

    kb = {
        "template_file": str(template_path.name),
        "slide_size_emu": {"width": prs.slide_width, "height": prs.slide_height},
        "layout_count": len(layouts),
        "slide_count": len(prs.slides),
        "semantic_layout_map": semantic_layout_map,
        "thanks_page_reference": {
            "expected_layout": "封底_End page",
            "expected_text": "Thanks.",
            "template_slide_index_0_based": 55,
        },
        "layouts": layouts,
        "sample_slides": sample_slides,
    }

    output_path.write_text(json.dumps(kb, ensure_ascii=False, indent=2), encoding="utf-8")
    return kb


def clear_all_slides(prs: Presentation) -> None:
    """Remove all existing slides from a presentation."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rel_id = sld_id.rId
        prs.part.drop_rel(rel_id)
        sld_id_lst.remove(sld_id)


def find_layout_by_name(prs: Presentation, keyword: str) -> object:
    """Find a slide layout using a partial name match."""
    for layout in prs.slide_layouts:
        if keyword in layout.name:
            return layout
    raise ValueError(f"Layout not found with keyword: {keyword}")


def fill_placeholders(slide: object, texts: list[str]) -> None:
    """Fill text placeholders in order of placeholder index."""
    placeholders = []
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        idx = shape.placeholder_format.idx
        placeholders.append((idx, shape))
    placeholders.sort(key=lambda item: item[0])

    for i, text in enumerate(texts):
        if i >= len(placeholders):
            break
        placeholders[i][1].text = text


def create_demo_ppt_from_template(template_path: Path, output_path: Path) -> None:
    """Reuse template layouts to generate a small demo deck."""
    prs = Presentation(str(template_path))
    clear_all_slides(prs)

    cover_layout = find_layout_by_name(prs, "封面_Cover")
    toc_layout = find_layout_by_name(prs, "目录页_Content")
    section_layout = find_layout_by_name(prs, "章节页_Section page")
    content_layout = find_layout_by_name(prs, "标准内容页（小标题）_Standard page with subtitle")

    cover = prs.slides.add_slide(cover_layout)
    fill_placeholders(cover, ["cc Demo Deck", "Generated from template layouts"])

    toc = prs.slides.add_slide(toc_layout)
    fill_placeholders(toc, ["目录", "项目背景\n方案设计\n实施计划"])

    section = prs.slides.add_slide(section_layout)
    fill_placeholders(section, ["项目背景 - Project Background", "本节介绍目标与范围", "01"])

    content = prs.slides.add_slide(content_layout)
    fill_placeholders(
        content,
        [
            "方案说明",
            "这里是随便内容：本页复用模板内容页样式，保持品牌视觉一致。",
            "副标题 - Subtitle",
        ],
    )

    end_layout = find_layout_by_name(prs, "封底_End page")
    end = prs.slides.add_slide(end_layout)
    fill_placeholders(end, ["Thanks."])

    prs.save(str(output_path))


def main() -> None:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    extracted = extract_media(TEMPLATE_PATH, ASSET_DIR)
    kb = build_template_knowledge_base(TEMPLATE_PATH, KB_PATH)

    create_demo_ppt_from_template(TEMPLATE_PATH, OUTPUT_PATH)
    print(f"Done. Generated: {OUTPUT_PATH.resolve()}")
    print(f"Extracted media count: {len(extracted)}")
    print(f"Knowledge base: {KB_PATH.resolve()}")
    print(f"Semantic end layout: {kb['semantic_layout_map']['end']}")


if __name__ == "__main__":
    main()
